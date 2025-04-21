import numpy as np
from faiss import IndexFlatL2
from transformers import AutoTokenizer, AutoModel, AutoModelForCausalLM
import torch
import os
import requests
import opik
from opik import track, opik_context
import streamlit as st
import requests
import fitz  
import pandas as pd  
from io import StringIO
import docx 
from pptx import Presentation  
import json
from firecrawl import FirecrawlApp
from IPython.display import Markdown, display
import time
import base64
from phi.agent import Agent, RunResponse
from phi.model.huggingface import HuggingFaceChat
from phi.utils.pprint import pprint_run_response
from phi.model.openai import OpenAIChat
from validator import LlmRagEvaluator, HallucinationPrompt, QACorrectnessPrompt
from guardrails import Guard
import phoenix.evals
from smolagents import CodeAgent, HfApiModel
from streamlit_pdf_viewer import pdf_viewer
from streamlit import session_state as ss
from phoenix.evals import (
    QA_PROMPT_RAILS_MAP,
    QA_PROMPT_TEMPLATE,
    OpenAIModel,
    llm_classify,
)
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")

# Initialize environment variable and model
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
ollama_url = "http://localhost:11434/api/generate"
model_path = os.getenv("MODEL_PATH") 
hf_model_path = os.getenv("HF_MODEL_PATH") 
tokenizer = AutoTokenizer.from_pretrained(model_path)
model = AutoModel.from_pretrained(model_path)
tokenizer_hf = AutoTokenizer.from_pretrained(hf_model_path)
model_hf = AutoModelForCausalLM.from_pretrained(hf_model_path)
model_hf.to(device)

# Function to generate embeddings
def get_embeddings(texts):
    inputs = tokenizer(texts, padding=True, truncation=True, return_tensors="pt").to(device)
    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1).cpu().numpy()
    return embeddings

# Function to generate embeddings and store them in Faiss
def generate_embeddings(documents, file_names):
    embeddings = get_embeddings(documents)
    index = IndexFlatL2(embeddings.shape[1])
    index.add(embeddings)
    return embeddings, index, file_names
    

# Function to find the most relevant document based on Faiss similarity search
def find_relevant_document(question, faiss_index, documents, file_names):
    question_embedding = get_embeddings([question])
    D, I = faiss_index.search(question_embedding, k = len(documents) )
    file_name = file_names[I[0][0]]
    return documents[I[0][0]], D[0][0], file_name  


@track(project_name="Document QA Agent", tags=['agent', 'python-library', 'querying'])
def generate_answer(agent, predefined_question, relevant_doc):
    prompt = f"Document: {relevant_doc}\n\nQuestion: {predefined_question}\nAnswer:"
    answer: RunResponse = agent.run(prompt)
    return answer

# Streamlit frontend
st.title('AI Ally: Document and Website QA')
st.write('Choose which insights do you want:')

# Create two columns for cards
col1, col2 = st.columns(2, gap="small", vertical_alignment="center")

# Card 1: Document Upload and QA
with col1:
    if st.button("Document Insights", use_container_width=True):
        # Hide other sections and show Document Upload section
        st.session_state["section"] = "document_upload"

# Card 2: Website Scraping and QA
with col2:
    if st.button("Website Insights", use_container_width=True):
        # Hide other sections and show Website Scraping section
        st.session_state["section"] = "website_scraping"

# Handle what section to display based on button clicked
if "section" not in st.session_state:
    st.session_state["section"] = "none"

# Section for Document Upload and QA
if st.session_state["section"] == "document_upload":
    uploaded_files = st.file_uploader("Choose files", type=["pdf", "txt", "docx", "xlsx", "pptx", "csv"], accept_multiple_files=True)

    # Function to display PDFs 
    def display_pdf(file):
        # Read file as bytes:
        bytes_data = file.getvalue()
        pdf_viewer(input=bytes_data, width=700)

    # Function to preview DOCX content
    def display_docx(file):
        doc = docx.Document(file)
        doc_text = ""
        for para in doc.paragraphs:
            doc_text += para.text + "\n"
        st.text_area("DOCX Content", doc_text, height=300)
    
    # Function to preview Excel content
    def display_excel(file):
        df = pd.read_excel(file)
        st.write(df)
    
    # Function to preview PPTX content
    def display_pptx(file):
        prs = Presentation(file)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"
        st.text_area("PPTX Content", pptx_text, height=300)

    # Helper function to read text from DOCX files
    def extract_text_from_docx(docx_file):
        doc = docx.Document(docx_file)
        doc_text = ""
        for para in doc.paragraphs:
            doc_text += para.text + "\n"
        return doc_text

    # Helper function to read text from Excel files
    def extract_text_from_excel(excel_file):
        df = pd.read_excel(excel_file)
        # Concatenate all columns into a single string
        excel_text = ""
        for col in df.columns:
            excel_text += "\n".join(df[col].astype(str)) + "\n"
        return excel_text

    # Helper function to read text from PPT files
    def extract_text_from_pptx(pptx_file):
        prs = Presentation(pptx_file)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"
        return pptx_text

    # Process each uploaded file
    if uploaded_files:
        all_documents = []
        file_names = []
        for uploaded_file in uploaded_files:
            # Process PDF files
            if uploaded_file.type == "application/pdf":
                pdf_file = fitz.open(stream=uploaded_file.read(), filetype="pdf")
                doc_text = ""
                for page in pdf_file:
                    doc_text += page.get_text()
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                preview_key = f"preview_{uploaded_file.name}"

                if st.button(f"Preview {uploaded_file.name}"):
                    if preview_key not in st.session_state:
                        st.session_state[preview_key] = False  # Default to not showing the preview
    
                    # Toggle the preview state
                    st.session_state[preview_key] = not st.session_state[preview_key]
    
                # If the preview is active, display the PDF
                if preview_key in st.session_state and st.session_state[preview_key]:
                    display_pdf(uploaded_file)
                     

            # Process DOCX files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc_text = extract_text_from_docx(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                preview_key = f"preview_{uploaded_file.name}"

                if st.button(f"Preview {uploaded_file.name}"):
                    if preview_key not in st.session_state:
                        st.session_state[preview_key] = False  # Default to not showing the preview
    
                    # Toggle the preview state
                    st.session_state[preview_key] = not st.session_state[preview_key]
    
                # If the preview is active, display the PDF
                if preview_key in st.session_state and st.session_state[preview_key]:
                    display_docx(uploaded_file)

            # Process TXT files
            elif uploaded_file.type == "text/plain":
                doc_text = StringIO(uploaded_file.getvalue().decode("utf-8")).read()
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                preview_key = f"preview_{uploaded_file.name}"

                if st.button(f"Preview {uploaded_file.name}"):
                    if preview_key not in st.session_state:
                        st.session_state[preview_key] = False  # Default to not showing the preview
    
                    # Toggle the preview state
                    st.session_state[preview_key] = not st.session_state[preview_key]
    
                # If the preview is active, display the PDF
                if preview_key in st.session_state and st.session_state[preview_key]:
                    st.text_area("Text File Content", doc_text, height=300)
            
            # Process Excel files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                doc_text = extract_text_from_excel(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                preview_key = f"preview_{uploaded_file.name}"

                if st.button(f"Preview {uploaded_file.name}"):
                    if preview_key not in st.session_state:
                        st.session_state[preview_key] = False  # Default to not showing the preview
    
                    # Toggle the preview state
                    st.session_state[preview_key] = not st.session_state[preview_key]
    
                # If the preview is active, display the PDF
                if preview_key in st.session_state and st.session_state[preview_key]:
                    display_excel(uploaded_file)
            
            # Process PPTX files
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                doc_text = extract_text_from_pptx(uploaded_file)
                all_documents.append(doc_text)
                file_names.append(uploaded_file.name)
                preview_key = f"preview_{uploaded_file.name}"

                if st.button(f"Preview {uploaded_file.name}"):
                    if preview_key not in st.session_state:
                        st.session_state[preview_key] = False  # Default to not showing the preview
    
                    # Toggle the preview state
                    st.session_state[preview_key] = not st.session_state[preview_key]
    
                # If the preview is active, display the PDF
                if preview_key in st.session_state and st.session_state[preview_key]:
                    display_pptx(uploaded_file)
            
            # Process CSV files
            elif uploaded_file.type == "text/csv":
                df = pd.read_csv(uploaded_file)
                csv_text = df.to_string(index=False)
                all_documents.append(csv_text)
                file_names.append(uploaded_file.name)
                preview_key = f"preview_{uploaded_file.name}"

                if st.button(f"Preview {uploaded_file.name}"):
                    if preview_key not in st.session_state:
                        st.session_state[preview_key] = False  # Default to not showing the preview
    
                    # Toggle the preview state
                    st.session_state[preview_key] = not st.session_state[preview_key]
    
                # If the preview is active, display the PDF
                if preview_key in st.session_state and st.session_state[preview_key]:
                    st.write(df)

        # After processing all files, generate embeddings for the combined documents
        embeddings, index, file_names = generate_embeddings(all_documents, file_names)
        
        agent = CodeAgent(tools = [], model=HfApiModel())
 
        example_questions = []
        for file in file_names:
            example_questions.extend([
                f"What is the main topic of {file}?",
                f"What are the key points of {file}?",
                f"Can you extract any conclusions from {file}?"
            ])
        
        # Display the predefined questions in a row
        st.subheader("Query:")
        
        # Ask a custom question
        user_question = st.text_input('Ask a question specifying the file name:')
        
        # Create columns dynamically for each question based on the number of predefined questions
        questions_per_row = 3  
        rows = (len(example_questions) + questions_per_row - 1) // questions_per_row  # Calculate the number of rows
        
        # Initialize a list to store the answers
        answers = []
        st.write('Example Questions:')
        # Display questions in rows 
        for i in range(rows):
            cols = st.columns(questions_per_row)  
            start_idx = i * questions_per_row
            end_idx = min((i + 1) * questions_per_row, len(example_questions))  # Ensure the last row doesn't exceed the list length
            
            # Add a button for each predefined question in the current row
            for j, predefined_question in enumerate(example_questions[start_idx:end_idx]):
                with cols[j]:
                    if st.button(predefined_question):
                        # Determine which file the question is related to based on its index
                        relevant_doc, distance, file_name = find_relevant_document(predefined_question, index, all_documents, file_names)
                        answer = generate_answer(agent, predefined_question, relevant_doc)
                        answers.append({
                            'question': predefined_question,
                            'file': file_name,
                            'distance': distance,
                            'answer': answer
                        })
        
        # Process the custom question if provided
        if user_question:
            relevant_doc, distance, file_name = find_relevant_document(user_question, index, all_documents, file_names)
            answer = generate_answer(agent, user_question, relevant_doc)
            answers.append({
                'question': user_question,
                'file': file_name,
                'distance': distance,
                'answer': answer
            })
        
        # Display all responses 
        if answers:
            st.subheader("Response:")
            for ans in answers:
                st.write(f"**Question:** {ans['question']}")
                st.write(f"**Found in file:** {ans['file']}")
                st.write(f"**Relevant document distance:** {ans['distance']}")
                st.write(f"**Answer:** {ans['answer']}")

# Section for Website Scraping and QA
if st.session_state["section"] == "website_scraping":
    url_input = st.text_input("Enter the URL of the website:")    
    if url_input:
        documents = []
        file_names = [url_input] 
        firecrawl_api_key = os.getenv('FIRECRAWL_API')

        # Initialize FireCrawlApp instance
        app = FirecrawlApp(api_key=firecrawl_api_key)

        try:
            # Scrape data using FireCrawl
            scrape_result = app.scrape_url(url_input, params={'formats': ['markdown', 'html']})

            # Extract content from the markdown field
            if 'markdown' in scrape_result:
                scraped_text = scrape_result['markdown']

                if scraped_text:
                    st.write(f"Preview website: {url_input}")
                    documents.append(scraped_text)
                else:
                    st.write(f"No valid content extracted for {url_input}. Please check the URL.")
            else:
                st.write(f"Error: 'markdown' key not found in the FireCrawl response. Full response shown above.")
        except Exception as e:
            st.error(f"Error while scraping: {e}")


        # Generate embeddings for the scraped content
        if documents:
            embeddings, index, file_names = generate_embeddings(documents, file_names)
            
            agent = CodeAgent(tools= [], model=HfApiModel())

            example_questions = [
                f"What is the main topic of this website?",
                f"Can you summarize this website?",
                f"What are the key points of this website?",
                f"What is the purpose of this website?",
                f"Can you extract any conclusions from this website?"
            ]

            st.subheader('Query:')
            # Ask the user for a question about the scraped content
            user_question = st.text_input('Ask a question:')
            questions_per_row = 3 
            rows = (len(example_questions) + questions_per_row - 1) // questions_per_row  # Calculate the number of rows
            
            # Initialize a list to store the answers
            answers = []
            st.write('Example Questions:') 
            # Display questions in rows
            for i in range(rows):
                cols = st.columns(questions_per_row)  # Create n columns per row
                start_idx = i * questions_per_row
                end_idx = min((i + 1) * questions_per_row, len(example_questions))  # Ensure the last row doesn't exceed the list length
                
                # Add a button for each predefined question in the current row
                for j, predefined_question in enumerate(example_questions[start_idx:end_idx]):
                    with cols[j]:
                        if st.button(predefined_question):
                            relevant_doc, distance, file_name = find_relevant_document(predefined_question, index, documents, file_names)
                            answer = generate_answer(agent, predefined_question, relevant_doc)
                            answers.append({
                                'question': predefined_question,
                                'file': file_name,
                                'distance': distance,
                                'answer': answer
                            })
            # Process the custom question if provided
            if user_question:
                relevant_doc, distance, file_name = find_relevant_document(user_question, index, documents, file_names)
                answer = generate_answer(agent, user_question, relevant_doc)
                answers.append({
                    'question': user_question,
                    'file': file_name,
                    'distance': distance,
                    'answer': answer
                })
            
            # Display responses
            if answers:
                st.subheader("Response:")
                for ans in answers:
                    st.write(f"**Question:** {ans['question']}")
                    st.write(f"**Answer:** {ans['answer']}")
        else:
                st.write("No valid scraped content to generate embeddings.")

